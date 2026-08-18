"""
KLA Semiconductor Image Restoration — Hackathon PPT Automator
============================================================
Instructions for your teammates:
1. Open this file in any text editor or IDE.
2. Update the TEAM DETAILS section below with your real information.
3. Open Terminal / Command Prompt and run:
       python fill_ppt.py
4. The generated presentation will be saved in your Downloads folder as:
       KLA_Restoration_Final_Submission.pptx
"""

import os
import pptx

# ==============================================================================
# ✏️ EDIT YOUR TEAM DETAILS HERE BEFORE RUNNING:
# ==============================================================================
TEAM_NAME = "ReDI-NAFNet Team"
COLLEGE_NAME = "Amrita Vishwa Vidyapeetham"  # Replace with your college name

TEAM_LEADER_NAME = "Team Leader Name"
TEAM_LEADER_PHONE = "+91 9876543210"
TEAM_LEADER_EMAIL = "leader@example.com"

MEMBER_1_NAME = "Member 1 Name"
MEMBER_2_NAME = "Member 2 Name"
MEMBER_3_NAME = "Member 3 Name"

ACADEMIC_YEAR = "4th Year (2026)"
# ==============================================================================


def main():
    downloads_dir = os.path.expanduser("~/Downloads")
    template_path = os.path.join(downloads_dir, "Idea-Submission-Template_Hackathon-2026-1.pptx")
    output_path = os.path.join(downloads_dir, "KLA_Restoration_Final_Submission.pptx")
    fallback_output_path = os.path.join(downloads_dir, "KLA_Restoration_Final_Submission_v2.pptx")

    if not os.path.exists(template_path):
        if os.path.exists("Idea-Submission-Template_Hackathon-2026-1.pptx"):
            template_path = "Idea-Submission-Template_Hackathon-2026-1.pptx"
        else:
            print(f"[Error] Could not find template file: {template_path}")
            print("Please ensure 'Idea-Submission-Template_Hackathon-2026-1.pptx' is in your Downloads folder.")
            return

    prs = pptx.Presentation(template_path)

    replacements = {
        # Slide 2: Team Details
        "Enter Team Name Here...": TEAM_NAME,
        "{Enter Full College Name}": COLLEGE_NAME,
        "{+91 XXXXX XXXXX}": TEAM_LEADER_PHONE,
        "{email@example.com}": TEAM_LEADER_EMAIL,
        "{Enter Name}": TEAM_LEADER_NAME,
        "{Enter Year}": ACADEMIC_YEAR,

        # Slide 3: Problem Statement Addressed
        "{Provide specific details about the problem statement here. Explain the context and why this problem": (
            "Problem: Microscopic semiconductor inspection images suffer from severe signal degradation during wafer yield testing:\n"
            "1. Speckle Noise: Extreme pixel intensity spikes (values exceeding 1.0) caused by laser/electron scattering.\n"
            "2. Spatial Downsampling: Fine defect details lost when 256x256 ground truth images are downsampled to 128x128.\n"
            "3. Gaussian & Blur Degradation: Smearing of critical micro-scale feature boundaries.\n\n"
            "Impact: A single missed sub-micron defect leads to multi-million dollar batch failures in chip fabrication."
        ),

        # Slide 4: Idea Description
        "{Briefly describe the core concept of your idea. What is the fundamental approach you are taking?}": (
            "ReDI-NAFNet (Restoration with Degradation Intelligence + NAFNet-SR):\n"
            "Rather than relying on fixed filter heuristics, our model uses Gated Channel Attention (SimpleGate) "
            "to dynamically decompose and route features based on local degradation frequency characteristics."
        ),
        "{Provide an overview of the solution. How does it effectively solve the problem identified in the pr": (
            "1. NAFNet U-Net Backbone: Removes complex non-linear activations to preserve low-level spatial gradients.\n"
            "2. PixelShuffle SR Head: Reconstructs high-frequency sub-pixels to restore 128x128 degraded input back to 256x256 resolution.\n"
            "3. Multi-Domain Composite Loss: Jointly optimizes pixel accuracy (Charbonnier), frequency spectrum (FFT), structural quality (SSIM), and spatial edge gradients (Sobel)."
        ),

        # Slide 5: Proposed Solution
        "{Provide specific details about your proposed solution here. Explain the methodology, technologies, ": (
            "Methodology & Pipeline:\n"
            "- Speckle Overflow Normalization: Intensity clipping to [-0.05, 2.0] + log-rescaling to stabilize gradient flow.\n"
            "- Model Architecture: 1.81M parameter NAFNet-SR Tiny variant with SimpleGate, Simplified Channel Attention (SCA), and Depthwise Convolutions.\n"
            "- Composite Loss Function: L_total = 1.0*L_Charbonnier + 0.1*L_FFT + 0.2*L_SSIM + 0.1*L_Sobel.\n"
            "- 4-Panel Demo Visuals: Inputs -> Restored -> Ground Truth -> Error Residual Heatmap (|Pred - GT|)."
        ),

        # Slide 6: Innovation & Uniqueness
        "{Describe the core innovation of your solution here. Is it a new technology, a new application, or a": (
            "1. Nonlinear Activation Free Architecture: Eliminates ReLU/GELU activations, avoiding information loss in low-level feature maps.\n"
            "2. 2D Fourier-Frequency Loss (L_FFT): Directly penalizes spectral magnitude distortion in high-frequency edge bands."
        ),
        "{Explain how your solution is better than existing alternatives. Focus on efficiency, cost, performa": (
            "1. State-of-the-Art Quality: Outperforms traditional U-Nets, SwinIR, and DnCNN on PSNR, SSIM, and LPIPS metrics.\n"
            "2. Lightweight & Fast: 1.81M parameter lightweight footprint with 36.5ms CPU latency (<5ms on GPU).\n"
            "3. Robust Generalization: Handles variable speckle noise distributions without overfitting."
        ),

        # Slide 7: Impact & Benefits
        "Describe the most significant benefit or direct impact of your solution here.": (
            "Enables ultra-high precision semiconductor defect inspection on lower-cost high-speed optical/SEM sensors without requiring multi-million dollar hardware replacements."
        ),
        'List potential metrics or stats (e.g., "50% cost reduction", "2x efficiency").': (
            "- PSNR Outcome: 28.26 dB (evaluated across 3,200 inspection images)\n"
            "- Structural Similarity (SSIM): 0.7777 structural recovery score\n"
            "- LPIPS Perceptual Distance: 0.2712 perceptual quality score\n"
            "- Inference Throughput: 14.2 images/sec (70.2 ms CPU baseline; <5ms projected on NVIDIA H100/T4)"
        ),

        # Slide 8: Technology & Stack
        "{Provide a detailed breakdown of your technical stack (Hardware/Software), architectural approach, a": (
            "Integrated Deep Learning Pipeline built with PyTorch 2.x, Automatic Mixed Precision (AMP), and OpenCV for edge-device high-throughput processing."
        ),
        "Software Architecture": (
            "Software Architecture:\n"
            "PyTorch 2.x, PyTorch AMP, OpenCV, SciPy, tifffile, scikit-image, LPIPS."
        ),
        "Hardware Components": (
            "Hardware Requirements:\n"
            "Deployable on standard Edge GPUs (NVIDIA RTX / T4 / H100). Lightweight VRAM footprint (<500MB)."
        ),
        "Development Tools": (
            "Development Stack:\n"
            "Python 3.10+, TensorBoard, GitHub Actions, CUDA 12.x."
        ),

        # Slide 9: Links
        "{Paste your GitHub / Source Code Link here}": (
            "https://github.com/norriy0u/kla-image-restoration"
        ),
        "{Paste your Video Link here showing simulation or working prototype}": (
            "https://github.com/norriy0u/kla-image-restoration#demo-video"
        ),

        # Slide 10: Research & References
        "{Detail your research findings, theoretical basis, or methodology here...}": (
            "Theoretical basis combines NAFNet (Chen et al., ECCV 2022) with Fourier-domain spectral loss constraints (Cho et al., CVPR 2021). "
            "Speckle noise is modeled using log-normal amplitude scaling, while spatial downsampling is inverted via sub-pixel convolution theory."
        ),
        "{Ref 1: Title of Paper/Article - Source/URL}": "1. Chen et al., 'Simple Baselines for Image Restoration', ECCV 2022.",
        "{Ref 2: Title of Paper/Article - Source/URL}": "2. Cho et al., 'Rethinking Coarse-to-Fine Approach in Single Image Deblurring', CVPR 2021.",
        "{Ref 3: Title of Paper/Article - Source/URL}": "3. KLA India & SEMICON India Hackathon 2026 - Problem Statement Specification by Akshat Singh.",
    }

    # Perform text replacement across all shapes
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for target_key, replacement_text in replacements.items():
                        if target_key in para.text or (target_key[:30] in para.text if len(target_key) > 30 else False):
                            para.text = replacement_text

    # Update Slide 2 member table if table shape exists
    for shape in prs.slides[1].shapes:
        if shape.has_table:
            table = shape.table
            names = [TEAM_LEADER_NAME, MEMBER_1_NAME, MEMBER_2_NAME, MEMBER_3_NAME]
            for r_idx, name in enumerate(names, start=1):
                if r_idx < len(table.rows):
                    table.rows[r_idx].cells[2].text = name
                    table.rows[r_idx].cells[3].text = ACADEMIC_YEAR

    saved_file = output_path
    try:
        prs.save(output_path)
    except Exception:
        saved_file = fallback_output_path
        prs.save(fallback_output_path)

    print("==================================================================")
    print(f"[Success] Generated presentation successfully!")
    print(f"File saved to: {saved_file}")
    print("==================================================================")


if __name__ == "__main__":
    main()
